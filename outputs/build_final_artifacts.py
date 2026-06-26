import json
import math
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path(r"C:\Users\jiten\btp_final2")
PROJECT = ROOT / "TSFFM-Depression-Detection"
RESULTS = PROJECT / "results"
OUT_DIR = ROOT / "outputs" / "final_artifacts"
ASSET_DIR = OUT_DIR / "assets"
REPORT_PATH = OUT_DIR / "TSFFM_Final_Evaluation_Report.docx"
PPTX_PATH = OUT_DIR / "TSFFM_Final_Evaluation_Presentation.pptx"


TEAM = "Jitendra Choudhary (22BEC059) and Kabir Singh Khair (22BEC061)"
SUPERVISOR = "Dr. Amit Vishwakarma"
INSTITUTE = "PDPM Indian Institute of Information Technology, Design and Manufacturing, Jabalpur"
DEPARTMENT = "Electronics and Communication Engineering"


PRIMARY = RGBColor(15, 37, 67)
ACCENT = RGBColor(21, 112, 110)
SUPPORT = RGBColor(184, 115, 51)
MUTED = RGBColor(91, 106, 125)
LIGHT = RGBColor(244, 247, 250)
TABLE_HEAD = "E8EEF5"
CALL_FILL = "F4F8F8"
RISK_FILL = "FFF4E8"


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def hex_to_rgb_tuple(hex_color):
    h = hex_color.strip("#")
    return tuple(int(h[i : i + 2], 16) / 255 for i in (0, 2, 4))


def read_metrics():
    main = {
        "accuracy": 57.67,
        "precision": 0.4688,
        "recall": 0.4255,
        "f1": 0.4461,
        "auc": 0.6597,
        "support": 352,
        "not_dep": 211,
        "dep": 141,
    }
    candidate = {
        "accuracy": 53.69,
        "precision": 0.4567,
        "recall": 0.8227,
        "f1": 0.5873,
        "auc": 0.5454,
    }
    hist_path = PROJECT / "backend" / "weights" / "training_history.json"
    history = {}
    if hist_path.exists():
        history = json.loads(hist_path.read_text(encoding="utf-8"))
    return main, candidate, history


def save_pipeline_diagram(path):
    fig, ax = plt.subplots(figsize=(12, 3.1), dpi=180)
    ax.axis("off")
    steps = [
        ("Video\nupload", "#0F2543"),
        ("Frame\nsampling\n5 FPS", "#15706E"),
        ("MediaPipe\nFaceMesh\n68 pts", "#B87333"),
        ("MediaPipe\nPose\n2 shoulders", "#7C5AA6"),
        ("Normalize\nand pad\n360 frames", "#4C6F91"),
        ("TSFFM\nBiLSTM\nattention", "#0F2543"),
        ("Prediction\nPDF report", "#15706E"),
    ]
    xs = np.linspace(0.05, 0.95, len(steps))
    y = 0.55
    for i, ((label, color), x) in enumerate(zip(steps, xs)):
        rect = plt.Rectangle((x - 0.055, y - 0.18), 0.11, 0.36, color=color, ec="none", alpha=0.96)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", color="white", fontsize=8.5, fontweight="bold")
        if i < len(steps) - 1:
            ax.annotate(
                "",
                xy=(xs[i + 1] - 0.065, y),
                xytext=(x + 0.065, y),
                arrowprops=dict(arrowstyle="->", color="#6B7280", lw=1.6),
            )
    ax.text(
        0.5,
        0.12,
        "End-to-end inference path used by the deployed prototype",
        ha="center",
        color="#334155",
        fontsize=10,
    )
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def save_model_diagram(path):
    fig, ax = plt.subplots(figsize=(11.5, 4.2), dpi=180)
    ax.axis("off")
    boxes = [
        (0.08, 0.68, "Face input\n360 x 272", "#E8EEF5", "#0F2543"),
        (0.28, 0.68, "Face stream\nMLP projection\n128-D", "#DDF3F2", "#15706E"),
        (0.08, 0.30, "Pose input\n360 x 8", "#F7EEE6", "#B87333"),
        (0.28, 0.30, "Body stream\nMLP projection\n32-D", "#F7EEE6", "#B87333"),
        (0.48, 0.50, "Fusion\n160-D sequence", "#EEF2FF", "#4C6F91"),
        (0.64, 0.50, "Temporal Conv1D\nfeature smoothing", "#E8EEF5", "#0F2543"),
        (0.79, 0.50, "BiLSTM\n2 x 128 hidden", "#DDF3F2", "#15706E"),
        (0.92, 0.50, "Attention + FC\n2 logits", "#F7EEE6", "#B87333"),
    ]
    for x, y, text, fill, edge in boxes:
        w, h = (0.14, 0.17) if x < 0.45 else (0.12, 0.18)
        rect = plt.Rectangle((x - w / 2, y - h / 2), w, h, color=fill, ec=edge, lw=1.6)
        ax.add_patch(rect)
        ax.text(x, y, text, ha="center", va="center", fontsize=8.2, color="#111827", fontweight="bold")
    arrows = [
        ((0.15, 0.68), (0.21, 0.68)),
        ((0.35, 0.68), (0.43, 0.55)),
        ((0.15, 0.30), (0.21, 0.30)),
        ((0.35, 0.30), (0.43, 0.45)),
        ((0.54, 0.50), (0.58, 0.50)),
        ((0.70, 0.50), (0.73, 0.50)),
        ((0.85, 0.50), (0.87, 0.50)),
    ]
    for start, end in arrows:
        ax.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", color="#475569", lw=1.5))
    ax.text(0.5, 0.12, "Final implementation: landmark-based two-stream fusion with temporal attention", ha="center", fontsize=10, color="#334155")
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def save_metrics_bar(path, main, candidate):
    labels = ["Accuracy", "Precision", "Recall", "F1", "AUC"]
    main_vals = [main["accuracy"] / 100, main["precision"], main["recall"], main["f1"], main["auc"]]
    cand_vals = [candidate["accuracy"] / 100, candidate["precision"], candidate["recall"], candidate["f1"], candidate["auc"]]
    x = np.arange(len(labels))
    width = 0.34
    fig, ax = plt.subplots(figsize=(9.5, 4.2), dpi=180)
    ax.bar(x - width / 2, main_vals, width, label="Selected checkpoint", color="#15706E")
    ax.bar(x + width / 2, cand_vals, width, label="Recall-oriented candidate", color="#B87333")
    ax.set_ylim(0, 1.0)
    ax.set_ylabel("Score")
    ax.set_title("Validation Performance Summary")
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, loc="upper left")
    for i, v in enumerate(main_vals):
        ax.text(i - width / 2, v + 0.02, f"{v:.2f}", ha="center", fontsize=8)
    for i, v in enumerate(cand_vals):
        ax.text(i + width / 2, v + 0.02, f"{v:.2f}", ha="center", fontsize=8)
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def save_workflow_diagram(path):
    fig, ax = plt.subplots(figsize=(10.5, 4.2), dpi=180)
    ax.axis("off")
    lanes = [
        ("React dashboard", "Upload form, patient/session metadata, progress states, result display", 0.80, "#0F2543"),
        ("FastAPI backend", "Validates video, saves temp file, orchestrates feature extraction and inference", 0.58, "#15706E"),
        ("ML services", "MediaPipe extractor, PyTorch model singleton, probability response", 0.36, "#B87333"),
        ("Report output", "ReportLab PDF with prediction, confidence, probabilities, and disclaimer", 0.14, "#4C6F91"),
    ]
    for title, desc, y, color in lanes:
        rect = plt.Rectangle((0.08, y - 0.07), 0.84, 0.13, color="#F8FAFC", ec=color, lw=1.8)
        ax.add_patch(rect)
        ax.text(0.11, y + 0.025, title, fontsize=10, fontweight="bold", color=color, ha="left", va="center")
        ax.text(0.11, y - 0.025, desc, fontsize=8.5, color="#334155", ha="left", va="center")
    for y1, y2 in [(0.73, 0.65), (0.51, 0.43), (0.29, 0.21)]:
        ax.annotate("", xy=(0.50, y2), xytext=(0.50, y1), arrowprops=dict(arrowstyle="->", color="#64748B", lw=1.4))
    ax.text(0.5, 0.04, "Prototype integration: model research connected to a usable screening interface", ha="center", fontsize=9.5, color="#334155")
    plt.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def generate_figures(main, candidate):
    figures = {
        "pipeline": ASSET_DIR / "system_pipeline.png",
        "model": ASSET_DIR / "model_architecture.png",
        "metrics": ASSET_DIR / "metrics_summary.png",
        "workflow": ASSET_DIR / "implementation_workflow.png",
    }
    save_pipeline_diagram(figures["pipeline"])
    save_model_diagram(figures["model"])
    save_metrics_bar(figures["metrics"], main, candidate)
    save_workflow_diagram(figures["workflow"])
    return figures


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_cell_text(cell, text, bold=False, color=None, size=9.5):
    cell.text = ""
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run(str(text))
    run.bold = bold
    run.font.size = Pt(size)
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def set_table_geometry(table, widths):
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = False
    for row in table.rows:
        for idx, width in enumerate(widths):
            row.cells[idx].width = Inches(width)


def set_run_font(run, name="Calibri", size=None, color=None, bold=None, italic=None):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:ascii"), name)
    run._element.rPr.rFonts.set(qn("w:hAnsi"), name)
    if size is not None:
        run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.bold = bold
    if italic is not None:
        run.italic = italic


def add_page_number(paragraph):
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = paragraph.add_run()
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_begin)
    run._r.append(instr)
    run._r.append(fld_end)


def configure_doc(doc):
    section = doc.sections[0]
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(0.78)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)
    section.header_distance = Inches(0.45)
    section.footer_distance = Inches(0.35)

    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Calibri"
    normal._element.rPr.rFonts.set(qn("w:ascii"), "Calibri")
    normal._element.rPr.rFonts.set(qn("w:hAnsi"), "Calibri")
    normal.font.size = Pt(10.8)
    normal.paragraph_format.space_after = Pt(5)
    normal.paragraph_format.line_spacing = 1.10

    for name, size, color, before, after in [
        ("Heading 1", 16, PRIMARY, 10, 6),
        ("Heading 2", 13, ACCENT, 8, 4),
        ("Heading 3", 11.5, RGBColor(31, 77, 120), 6, 3),
    ]:
        style = styles[name]
        style.font.name = "Calibri"
        style._element.rPr.rFonts.set(qn("w:ascii"), "Calibri")
        style._element.rPr.rFonts.set(qn("w:hAnsi"), "Calibri")
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = color
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(after)
        style.paragraph_format.keep_with_next = True

    if "Caption" in styles:
        cap = styles["Caption"]
        cap.font.name = "Calibri"
        cap._element.rPr.rFonts.set(qn("w:ascii"), "Calibri")
        cap._element.rPr.rFonts.set(qn("w:hAnsi"), "Calibri")
        cap.font.size = Pt(8.5)
        cap.font.italic = True
        cap.font.color.rgb = MUTED
        cap.paragraph_format.space_before = Pt(2)
        cap.paragraph_format.space_after = Pt(8)

    header = section.header.paragraphs[0]
    header.text = "TSFFM Depression Detection | Final Evaluation Report"
    header.alignment = WD_ALIGN_PARAGRAPH.LEFT
    set_run_font(header.runs[0], size=8.5, color=MUTED)

    footer = section.footer.paragraphs[0]
    r = footer.add_run("Page ")
    set_run_font(r, size=8.5, color=MUTED)
    add_page_number(footer)


def page_break(doc):
    doc.add_page_break()


def add_title_para(doc, text, size=24, color=PRIMARY, align=WD_ALIGN_PARAGRAPH.CENTER, bold=True, after=8):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_after = Pt(after)
    run = p.add_run(text)
    set_run_font(run, size=size, color=color, bold=bold)
    return p


def add_body(doc, text, bold_lead=None):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.space_after = Pt(5)
    if bold_lead and text.startswith(bold_lead):
        r1 = p.add_run(bold_lead)
        set_run_font(r1, bold=True)
        r2 = p.add_run(text[len(bold_lead) :])
        set_run_font(r2)
    else:
        r = p.add_run(text)
        set_run_font(r)
    return p


def add_bullets(doc, items, level=0):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.left_indent = Inches(0.28 + level * 0.22)
        p.paragraph_format.first_line_indent = Inches(-0.12)
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(item)
        set_run_font(r, size=10.2)


def add_numbered(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Number")
        p.paragraph_format.left_indent = Inches(0.3)
        p.paragraph_format.first_line_indent = Inches(-0.12)
        p.paragraph_format.space_after = Pt(3)
        r = p.add_run(item)
        set_run_font(r, size=10.2)


def add_callout(doc, title, text, fill=CALL_FILL):
    table = doc.add_table(rows=1, cols=1)
    table.autofit = False
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0)
    cell.width = Inches(6.25)
    set_cell_shading(cell, fill)
    p = cell.paragraphs[0]
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run(title)
    set_run_font(r, size=10.2, color=ACCENT, bold=True)
    p2 = cell.add_paragraph()
    p2.paragraph_format.space_after = Pt(0)
    r2 = p2.add_run(text)
    set_run_font(r2, size=9.5, color=RGBColor(44, 62, 80))
    doc.add_paragraph()


def add_caption(doc, text):
    p = doc.add_paragraph(style="Caption")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(text)


def add_picture(doc, path, width=6.0, caption=None):
    if Path(path).exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run()
        r.add_picture(str(path), width=Inches(width))
        if caption:
            add_caption(doc, caption)


def add_kv_table(doc, rows, widths=(1.8, 4.5), header=None):
    table = doc.add_table(rows=0, cols=2)
    table.style = "Table Grid"
    if header:
        cells = table.add_row().cells
        set_cell_shading(cells[0], TABLE_HEAD)
        set_cell_shading(cells[1], TABLE_HEAD)
        set_cell_text(cells[0], header[0], bold=True, color=PRIMARY)
        set_cell_text(cells[1], header[1], bold=True, color=PRIMARY)
    for key, value in rows:
        cells = table.add_row().cells
        set_cell_text(cells[0], key, bold=True, color=PRIMARY)
        set_cell_text(cells[1], value, size=9.5)
    set_table_geometry(table, widths)
    doc.add_paragraph()
    return table


def add_matrix_table(doc, headers, rows, widths):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    for i, h in enumerate(headers):
        set_cell_shading(table.rows[0].cells[i], TABLE_HEAD)
        set_cell_text(table.rows[0].cells[i], h, bold=True, color=PRIMARY, size=9)
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_text(cells[i], value, size=8.8)
    set_table_geometry(table, widths)
    doc.add_paragraph()
    return table


def add_section_title(doc, idx, title):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(2)
    r = p.add_run(f"{idx:02d}  ")
    set_run_font(r, size=9, color=SUPPORT, bold=True)
    r2 = p.add_run(title)
    set_run_font(r2, size=16, color=PRIMARY, bold=True)
    return p


def build_report(main, candidate, history, figures):
    doc = Document()
    configure_doc(doc)

    # Page 1: cover
    for _ in range(5):
        doc.add_paragraph()
    add_title_para(doc, "Depression Detection using Two-Stream Feature Fusion Model", size=24, after=2)
    add_title_para(doc, "Final Evaluation Report", size=20, color=ACCENT, after=12)
    add_title_para(doc, "Bachelor Thesis Project", size=13, color=MUTED, bold=True, after=18)
    rows = [
        ("Submitted by", TEAM),
        ("Supervisor", SUPERVISOR),
        ("Department", DEPARTMENT),
        ("Institute", INSTITUTE),
        ("Evaluation stage", "End-semester final evaluation"),
        ("Project stack", "PyTorch, MediaPipe, OpenCV, FastAPI, React/Vite, ReportLab"),
    ]
    add_kv_table(doc, rows, widths=(1.6, 4.7))
    add_callout(
        doc,
        "Medical disclaimer",
        "This report describes an AI-assisted screening prototype. It is not a clinical diagnostic system and must be validated by qualified mental-health professionals before any real clinical use.",
        fill=RISK_FILL,
    )

    # Page 2
    page_break(doc)
    add_section_title(doc, 1, "Abstract")
    add_body(
        doc,
        "This project develops an AI-assisted depression screening prototype that analyzes video interviews using facial and body-behavior features. The final system extends the mid-semester concept into an integrated pipeline with feature extraction, temporal modeling, backend inference, a web dashboard, probability visualization, and automatic PDF report generation."
    )
    add_body(
        doc,
        "The final implementation uses MediaPipe FaceMesh to extract a 68-landmark facial representation and MediaPipe Pose to capture shoulder-level body cues. Each video is sampled to a fixed 360-frame sequence, normalized for translation and scale, projected through separate face and body streams, fused into a 160-dimensional temporal sequence, and classified with a bidirectional LSTM plus temporal attention."
    )
    add_body(
        doc,
        "Validation on 352 segments shows 57.67 percent accuracy, 0.4688 precision, 0.4255 recall, 0.4461 F1-score, and 0.6597 AUC-ROC for the selected checkpoint. A recall-oriented candidate reaches 0.8227 recall and 0.5873 F1-score but lowers accuracy and AUC. These results demonstrate a functioning end-to-end prototype while also showing the need for larger data, stronger generalization, and clinical validation."
    )
    add_kv_table(
        doc,
        [
            ("Keywords", "Depression detection, TSFFM, facial landmarks, pose landmarks, BiLSTM, temporal attention, FastAPI, React"),
            ("Primary output", "Binary screening prediction: depressed or not depressed"),
            ("Evaluation basis", "Validation split metrics, confusion matrix, ROC curve, training history, and system integration tests"),
        ],
        widths=(1.6, 4.7),
    )

    # Page 3
    page_break(doc)
    add_section_title(doc, 2, "Table of Contents")
    toc_rows = [
        ("1", "Abstract", "2"),
        ("2", "Table of Contents", "3"),
        ("3", "List of Figures, Tables, and Abbreviations", "4"),
        ("4", "Introduction", "5"),
        ("5", "Problem Statement and Motivation", "6"),
        ("6", "Objectives and Scope", "7"),
        ("7", "Literature Survey", "8"),
        ("8", "Proposed System Overview", "9"),
        ("9", "Dataset and Label Design", "10"),
        ("10", "Preprocessing and Feature Extraction", "11"),
        ("11", "Model Architecture", "12"),
        ("12", "Training Strategy", "13"),
        ("13", "Evaluation Protocol", "14"),
        ("14", "Final Results", "15"),
        ("15", "Confusion Matrix and ROC Analysis", "16"),
        ("16", "Training Dynamics", "17"),
        ("17", "Discussion and Error Analysis", "18"),
        ("18", "Backend and API Implementation", "19"),
        ("19", "Frontend and Report Generation", "20"),
        ("20", "Testing and Deployment Readiness", "21"),
        ("21", "Ethics, Privacy, and Clinical Safety", "22"),
        ("22", "Limitations and Risk Mitigation", "23"),
        ("23", "Future Scope", "24"),
        ("24", "Conclusion and References", "25"),
    ]
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    for i, h in enumerate(["No.", "Section", "Page"]):
        set_cell_shading(table.rows[0].cells[i], TABLE_HEAD)
        set_cell_text(table.rows[0].cells[i], h, bold=True, color=PRIMARY)
    for row in toc_rows:
        cells = table.add_row().cells
        for i, val in enumerate(row):
            set_cell_text(cells[i], val, size=9.2)
    set_table_geometry(table, (0.7, 4.8, 0.8))

    # Page 4
    page_break(doc)
    add_section_title(doc, 3, "List of Figures, Tables, and Abbreviations")
    add_matrix_table(
        doc,
        ["Figure", "Description"],
        [
            ("Figure 1", "Final end-to-end inference pipeline"),
            ("Figure 2", "Implemented TSFFM-BiLSTM-attention architecture"),
            ("Figure 3", "Validation metric comparison"),
            ("Figure 4", "Confusion matrix and ROC curve"),
            ("Figure 5", "Training and validation history"),
            ("Figure 6", "Integrated web/backend/reporting workflow"),
        ],
        (1.2, 5.1),
    )
    add_matrix_table(
        doc,
        ["Abbreviation", "Meaning"],
        [
            ("TSFFM", "Two-Stream Feature Fusion Model"),
            ("BiLSTM", "Bidirectional Long Short-Term Memory network"),
            ("PHQ-8", "Eight-item Patient Health Questionnaire"),
            ("AUC-ROC", "Area under the Receiver Operating Characteristic curve"),
            ("API", "Application Programming Interface"),
            ("ROI", "Region of Interest"),
        ],
        (1.6, 4.7),
    )

    # Page 5
    page_break(doc)
    add_section_title(doc, 4, "Introduction")
    add_body(
        doc,
        "Depression is a major public-health challenge and can remain undetected when access to mental-health professionals is limited or when patients under-report symptoms. The World Health Organization currently estimates that approximately 332 million people worldwide have depression, which reinforces the need for scalable, accessible, and early screening support."
    )
    add_body(
        doc,
        "Clinical screening normally depends on interviews, questionnaires, and professional interpretation. These remain essential, but they also require time, trained personnel, and patient willingness to disclose internal states. Video-based behavioral analysis can provide a complementary signal by observing non-verbal cues such as facial expressiveness, posture, gaze behavior, and psychomotor activity."
    )
    add_body(
        doc,
        "The project therefore explores whether a two-stream visual model can learn discriminative behavioral patterns from interview videos. The final result is not merely a model script: it is a working prototype that connects preprocessing, model inference, web interaction, visualized probabilities, and PDF report generation into one demonstrable system."
    )
    add_callout(
        doc,
        "Final evaluation position",
        "The system should be judged as an AI-assisted screening prototype with research and engineering value, not as a deployable clinical diagnostic product.",
    )

    # Page 6
    page_break(doc)
    add_section_title(doc, 5, "Problem Statement and Motivation")
    add_body(
        doc,
        "The central problem is to automatically classify whether a subject shows depression-related indicators from a video interview. The system must process raw or preprocessed visual evidence, extract useful behavioral features, fuse face and body information, and return a binary prediction with confidence values that can be understood by a user."
    )
    add_bullets(
        doc,
        [
            "Traditional diagnosis can be subjective, time-intensive, and difficult to access in rural or resource-constrained settings.",
            "Self-reporting can be affected by stigma, recall bias, and under-reporting.",
            "Single-modality AI systems may miss important context when they observe only face, speech, posture, or text.",
            "Video is non-invasive and uses commonly available hardware, making it practical for a prototype screening workflow.",
        ],
    )
    add_body(
        doc,
        "The motivation carried from the mid-semester work was to combine facial expression and posture cues into a richer representation. The final work keeps this idea but upgrades the technical path from frame-level fusion toward sequence modeling, attention, and a full-stack user-facing implementation."
    )

    # Page 7
    page_break(doc)
    add_section_title(doc, 6, "Objectives and Scope")
    add_numbered(
        doc,
        [
            "Build a video-based depression screening pipeline using visual behavioral features.",
            "Extract facial and body features in a repeatable format suitable for sequence modeling.",
            "Design a two-stream fusion model that keeps facial and pose information separate before fusion.",
            "Add temporal learning so the model can use behavior over time rather than isolated frames.",
            "Evaluate the system using accuracy, precision, recall, F1-score, AUC-ROC, confusion matrix, and training curves.",
            "Expose the trained model through a FastAPI backend and an interactive React dashboard.",
            "Generate a structured PDF report for each analyzed video while preserving a clear clinical disclaimer.",
        ],
    )
    add_callout(
        doc,
        "Scope boundary",
        "The prototype performs binary screening using visual cues only. It does not include audio, transcript semantics, psychiatrist-facing calibration, or external clinical validation.",
    )

    # Page 8
    page_break(doc)
    add_section_title(doc, 7, "Literature Survey")
    add_matrix_table(
        doc,
        ["Approach", "Typical signal", "Strength", "Limitation"],
        [
            ("Speech analysis", "Pitch, energy, pauses, voice quality", "Captures affective and psychomotor cues", "Sensitive to microphone quality and language"),
            ("Facial expression", "Action units, gaze, landmark motion", "Directly observes affect display", "Can miss posture and whole-body behavior"),
            ("Body posture", "Skeleton joints, movement amplitude", "Observes energy, slumping, and motion", "May miss subtle facial emotion"),
            ("Text or transcript", "Sentiment, topics, lexical patterns", "Strong semantic signal", "Depends on speech content and transcription quality"),
            ("Multimodal fusion", "Two or more behavioral channels", "More complete evidence", "More synchronization and data complexity"),
        ],
        (1.25, 1.65, 1.85, 1.55),
    )
    add_body(
        doc,
        "The proposed TSFFM approach is positioned as a visual multimodal method. It avoids specialized hardware and focuses on face plus body cues, making it more practical than physiological sensing while being richer than a face-only pipeline."
    )

    # Page 9
    page_break(doc)
    add_section_title(doc, 8, "Proposed System Overview")
    add_picture(doc, figures["pipeline"], width=6.1, caption="Figure 1: Final end-to-end inference pipeline.")
    add_body(
        doc,
        "The final system accepts a video file, samples frames at 5 FPS, extracts MediaPipe face and pose landmarks, normalizes the sequence, and pads or truncates it to 360 frames. The model then projects face and body features through separate streams, fuses them, models temporal dynamics, and emits class probabilities."
    )
    add_body(
        doc,
        "The deployed backend returns prediction, confidence, class probabilities, detection rates, frame counts, processing time, and a report URL. This transforms the model from an isolated notebook-style artifact into a usable screening workflow."
    )

    # Page 10
    page_break(doc)
    add_section_title(doc, 9, "Dataset and Label Design")
    add_body(
        doc,
        "The project is organized around the DAIC-WOZ / E-DAIC style interview dataset format, where participants are recorded in semi-clinical interviews and depression labels are derived from PHQ-style scores. The training loader reads participant directories containing pre-extracted face keypoints, pose confidence features, and binary labels."
    )
    add_kv_table(
        doc,
        [
            ("Training input", "Pre-extracted NumPy arrays for face keypoints and pose confidence features"),
            ("Segment structure", "Each participant can contribute multiple temporal segments"),
            ("Label type", "Binary label: 0 for not depressed, 1 for depressed"),
            ("Validation support", "352 segments in the reported evaluation: 211 not depressed and 141 depressed"),
            ("Sequence design", "Each segment is downsampled to 360 frames for model input"),
        ],
        widths=(1.75, 4.55),
    )
    add_body(
        doc,
        "The validation distribution is moderately imbalanced, with more not-depressed samples than depressed samples. The training code addresses this using class weights, focal loss, and a weighted random sampler."
    )

    # Page 11
    page_break(doc)
    add_section_title(doc, 10, "Preprocessing and Feature Extraction")
    add_body(
        doc,
        "For deployed inference, OpenCV reads the uploaded video and MediaPipe extracts visual landmarks. FaceMesh is reduced to a stable 68-landmark subset, with each landmark represented by x, y, z, and confidence. Pose extraction currently uses the left and right shoulders, each represented by x, y, z, and visibility."
    )
    add_matrix_table(
        doc,
        ["Feature family", "Raw shape", "Flattened model input", "Purpose"],
        [
            ("Face landmarks", "360 x 68 x 4", "360 x 272", "Facial affect, gaze, mouth/eye movement proxies"),
            ("Pose landmarks", "360 x 2 x 4", "360 x 8", "Shoulder posture and upper-body orientation proxies"),
            ("Metadata", "Detection rates, FPS, sampled frames", "API response fields", "Quality visibility for users"),
        ],
        (1.4, 1.45, 1.55, 1.9),
    )
    add_body(
        doc,
        "The normalization routine subtracts the per-frame mean and divides by the maximum distance from the center. This improves translation and scale invariance so that absolute camera position contributes less to the classifier."
    )

    # Page 12
    page_break(doc)
    add_section_title(doc, 11, "Model Architecture")
    add_picture(doc, figures["model"], width=6.15, caption="Figure 2: Implemented TSFFM-BiLSTM-attention architecture.")
    add_body(
        doc,
        "The implemented model differs from the earlier mid-semester diagram. Instead of using full image CNN streams, the final code uses compact landmark-feature projections. Face features are projected from 272 dimensions to 128 dimensions and pose features from 8 dimensions to 32 dimensions. The streams are concatenated into a 160-dimensional temporal sequence."
    )
    add_body(
        doc,
        "A one-dimensional convolution smooths the fused temporal sequence before it enters a bidirectional LSTM. Temporal attention then learns which frames are most informative before the final fully connected classifier outputs two logits."
    )

    # Page 13
    page_break(doc)
    add_section_title(doc, 12, "Training Strategy")
    add_body(
        doc,
        "The training script uses PyTorch with Adam optimization, weight decay, class weighting, focal loss, gradient clipping, balanced sampling, learning-rate reduction on plateau, and early stopping based on validation F1. These choices are appropriate for a small, imbalanced behavioral dataset where depressed samples may be harder to classify reliably."
    )
    add_matrix_table(
        doc,
        ["Training element", "Final setting or behavior"],
        [
            ("Optimizer", "Adam with weight decay"),
            ("Loss", "Focal loss with class weights"),
            ("Sampling", "WeightedRandomSampler for class-balanced batches"),
            ("Sequence augmentation", "Temporal jittering and small Gaussian landmark noise during training"),
            ("Checkpointing", "Best model saved using validation metric improvement"),
            ("Regularization", "Dropout, weight decay, gradient clipping, learning-rate scheduling"),
        ],
        (2.0, 4.3),
    )
    add_body(
        doc,
        "The final training history contains 19 logged epochs. Training accuracy increases strongly, while validation accuracy and AUC fluctuate, indicating that generalization remains the main research challenge."
    )

    # Page 14
    page_break(doc)
    add_section_title(doc, 13, "Evaluation Protocol")
    add_body(
        doc,
        "The evaluation script runs the trained model on the validation split and records class predictions, class probabilities, and ground-truth labels. It computes standard binary-classification metrics and saves a text report, confusion matrix, ROC curve, and training-history plots."
    )
    add_kv_table(
        doc,
        [
            ("Accuracy", "Overall fraction of correct predictions"),
            ("Precision", "Among predicted depressed cases, how many were actually depressed"),
            ("Recall", "Among actual depressed cases, how many were detected"),
            ("F1-score", "Harmonic mean of precision and recall"),
            ("AUC-ROC", "Ranking quality across probability thresholds"),
            ("Confusion matrix", "Class-wise error pattern for false positives and false negatives"),
        ],
        widths=(1.6, 4.7),
    )
    add_callout(
        doc,
        "Clinical screening interpretation",
        "Recall is especially important for a screening assistant because missed depressed cases are high-risk. However, high recall alone is insufficient if false positives become excessive.",
    )

    # Page 15
    page_break(doc)
    add_section_title(doc, 14, "Final Results")
    add_picture(doc, figures["metrics"], width=6.1, caption="Figure 3: Validation metric comparison.")
    add_matrix_table(
        doc,
        ["Checkpoint", "Accuracy", "Precision", "Recall", "F1", "AUC"],
        [
            ("Selected checkpoint", "57.67%", "0.4688", "0.4255", "0.4461", "0.6597"),
            ("Recall-oriented candidate", "53.69%", "0.4567", "0.8227", "0.5873", "0.5454"),
        ],
        (1.65, 0.9, 0.9, 0.9, 0.9, 0.9),
    )
    add_body(
        doc,
        "The selected checkpoint has the better AUC and accuracy, while the recall-oriented candidate is more aggressive in detecting depressed cases. This trade-off is important for final evaluation because it shows the model can be tuned toward sensitivity but requires threshold calibration and better validation before practical use."
    )

    # Page 16
    page_break(doc)
    add_section_title(doc, 15, "Confusion Matrix and ROC Analysis")
    row_table = doc.add_table(rows=1, cols=2)
    row_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, img in enumerate([RESULTS / "confusion_matrix.png", RESULTS / "roc_curve.png"]):
        cell = row_table.rows[0].cells[i]
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if img.exists():
            p.add_run().add_picture(str(img), width=Inches(3.0))
    doc.add_paragraph(style="Caption").add_run("Figure 4: Confusion matrix and ROC curve from the selected evaluation.")
    add_body(
        doc,
        "The selected model correctly classifies more not-depressed samples than depressed samples. The AUC-ROC of 0.6597 suggests the probability scores contain useful ranking information, but the operating threshold still needs calibration to reach clinically acceptable sensitivity."
    )
    add_body(
        doc,
        "Approximate confusion-matrix interpretation from the report is 143 true negatives, 68 false positives, 81 false negatives, and 60 true positives. The false-negative count is the most important weakness for a mental-health screening context."
    )

    # Page 17
    page_break(doc)
    add_section_title(doc, 16, "Training Dynamics")
    row_table = doc.add_table(rows=1, cols=2)
    row_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, img in enumerate([RESULTS / "accuracy_history.png", RESULTS / "loss_history.png"]):
        cell = row_table.rows[0].cells[i]
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if img.exists():
            p.add_run().add_picture(str(img), width=Inches(3.0))
    doc.add_paragraph(style="Caption").add_run("Figure 5: Training and validation history.")
    if history:
        final_train = history.get("train_acc", [None])[-1]
        best_val = max(history.get("val_acc", [0]))
        best_auc = max(history.get("val_auc", [0]))
        add_kv_table(
            doc,
            [
                ("Final logged training accuracy", f"{final_train:.2f}%" if final_train is not None else "N/A"),
                ("Best logged validation accuracy", f"{best_val:.2f}%"),
                ("Best logged validation AUC", f"{best_auc:.4f}"),
            ],
            widths=(2.3, 4.0),
        )
    add_body(
        doc,
        "The widening difference between training and validation behavior suggests overfitting. This is expected for small behavioral datasets and should guide the next iteration toward stronger regularization, participant-level split checks, and more data diversity."
    )

    # Page 18
    page_break(doc)
    add_section_title(doc, 17, "Discussion and Error Analysis")
    add_body(
        doc,
        "The end-semester prototype demonstrates a complete technical path, but the current results should be interpreted cautiously. Depression-related visual behavior is subtle, heterogeneous, and influenced by culture, personality, interview setup, camera position, lighting, and co-occurring conditions."
    )
    add_bullets(
        doc,
        [
            "The selected checkpoint is conservative and misses many depressed samples.",
            "The recall-oriented candidate detects many more depressed samples but increases false positives.",
            "The deploy extractor uses a limited pose representation, so body-language evidence is under-used compared with the original design goal.",
            "Validation performance indicates the model has learned some discriminative signal but has not yet reached clinical-grade robustness.",
        ],
    )
    add_callout(
        doc,
        "Most important final insight",
        "The project succeeds as an integrated proof of concept. The main remaining research task is not interface engineering but improving generalization and clinically meaningful sensitivity.",
    )

    # Page 19
    page_break(doc)
    add_section_title(doc, 18, "Backend and API Implementation")
    add_picture(doc, figures["workflow"], width=6.1, caption="Figure 6: Integrated web/backend/reporting workflow.")
    add_body(
        doc,
        "The FastAPI backend exposes a `/api/predict` endpoint. It validates the uploaded file extension, saves the temporary video, extracts features, runs model inference, generates a PDF report, returns structured JSON, and cleans up temporary uploads."
    )
    add_matrix_table(
        doc,
        ["Response field", "Meaning"],
        [
            ("prediction", "depressed or not_depressed"),
            ("confidence", "Probability of the predicted class"),
            ("depression_probability", "Softmax probability for depressed class"),
            ("face/body detection rate", "Extractor quality indicators"),
            ("frames_processed", "Effective sequence evidence used"),
            ("report_url", "Download path for the generated PDF report"),
        ],
        (2.0, 4.3),
    )

    # Page 20
    page_break(doc)
    add_section_title(doc, 19, "Frontend and Report Generation")
    add_body(
        doc,
        "The React/Vite frontend turns the model into an accessible screening dashboard. It supports drag-and-drop video upload, patient reference entry, session-date entry, progress stages, result display, probability visualization, and report download."
    )
    add_bullets(
        doc,
        [
            "UploadBox validates video format and collects patient/session metadata.",
            "App.jsx coordinates API submission, loading states, result states, and retry behavior.",
            "ResultCard presents positive or negative screening with confidence and clinical recommendation text.",
            "ProbabilityChart visualizes depressed, not-depressed, face-detection, and body-detection confidence values.",
            "ReportLab generates a PDF screening report with metadata, probabilities, technical details, and a disclaimer.",
        ],
    )
    add_body(
        doc,
        "This engineering layer is a major end-semester improvement over the mid-semester model-only demonstration, because it shows how a trained model can be packaged for real user interaction."
    )

    # Page 21
    page_break(doc)
    add_section_title(doc, 20, "Testing and Deployment Readiness")
    add_body(
        doc,
        "The repository contains generated PDF reports, uploaded demo videos, built frontend assets, backend logs, and evaluation artifacts. These indicate that the major system paths have been exercised locally."
    )
    add_matrix_table(
        doc,
        ["Area", "Current readiness", "Recommended final check"],
        [
            ("Backend", "FastAPI app and prediction route implemented", "Run local server and verify `/docs` plus `/api/predict`"),
            ("Model", "Weights and history available", "Confirm model-layer consistency across train/evaluate/inference scripts"),
            ("Frontend", "Vite app and production dist present", "Run upload workflow with at least one known demo video"),
            ("Reports", "PDF reports generated in backend/reports", "Check wording, disclaimer, and file cleanup behavior"),
            ("Packaging", "Local project structure complete", "Add Docker or documented deployment script"),
        ],
        (1.25, 2.25, 2.8),
    )
    add_body(
        doc,
        "Before a final demo, the most important technical check is to run one end-to-end upload from the frontend and confirm that the generated probabilities, PDF report, and UI result agree."
    )

    # Page 22
    page_break(doc)
    add_section_title(doc, 21, "Ethics, Privacy, and Clinical Safety")
    add_body(
        doc,
        "Mental-health AI systems require careful ethical boundaries. Video interviews are sensitive biometric and behavioral data. Any real deployment would require informed consent, secure storage, access controls, retention limits, and review by clinical and institutional ethics bodies."
    )
    add_bullets(
        doc,
        [
            "The interface and report must state clearly that the result is not a diagnosis.",
            "The system should not be used for employment, insurance, disciplinary, or high-stakes decisions.",
            "Bias testing is required across gender, age, language, skin tone, camera quality, and cultural expression patterns.",
            "Raw uploaded videos should be deleted after processing unless explicit consent and governance exist.",
            "Human clinical review must remain mandatory for any positive or negative screening result.",
        ],
    )
    add_callout(
        doc,
        "Safety principle",
        "The safest role for the prototype is assistive triage and research exploration, not autonomous mental-health judgment.",
        fill=RISK_FILL,
    )

    # Page 23
    page_break(doc)
    add_section_title(doc, 22, "Limitations and Risk Mitigation")
    add_matrix_table(
        doc,
        ["Limitation", "Impact", "Mitigation"],
        [
            ("Small clinical datasets", "Overfitting and weak generalization", "Use larger, diverse, participant-separated datasets"),
            ("Visual-only evidence", "Misses speech and language symptoms", "Add audio and transcript streams"),
            ("Limited deployed pose points", "Body stream under-represents movement", "Use more joints and temporal motion descriptors"),
            ("Threshold uncertainty", "Recall/precision trade-off remains unstable", "Calibrate thresholds with clinician-guided validation"),
            ("Dataset bias", "Potential demographic unfairness", "Run subgroup evaluation and fairness audits"),
        ],
        (1.55, 2.0, 2.75),
    )
    add_body(
        doc,
        "One implementation caveat is model-version consistency. The inference service loads the final two-layer BiLSTM configuration, while evaluation code should be reviewed to ensure it instantiates the same architecture used by the saved checkpoint. This should be aligned before final submission or deployment."
    )

    # Page 24
    page_break(doc)
    add_section_title(doc, 23, "Future Scope")
    add_numbered(
        doc,
        [
            "Add an audio stream using pitch, energy, pause duration, speech rate, and voice-quality descriptors.",
            "Add a transcript/text stream using clinically meaningful language features and modern contextual embeddings.",
            "Use participant-level cross-validation and external validation on additional datasets.",
            "Expand pose features from shoulders to full upper-body and motion descriptors.",
            "Try temporal transformers or multimodal attention for richer sequence fusion.",
            "Calibrate probabilities and operating thresholds for screening sensitivity.",
            "Package the backend and frontend with Docker and add a deployment guide.",
            "Conduct clinical review with mental-health professionals before any real-world evaluation.",
        ],
    )
    add_body(
        doc,
        "The strongest next step is a multimodal tri-stream model that combines visual behavior, vocal behavior, and transcript semantics while keeping privacy and clinical validation at the center."
    )

    # Page 25
    page_break(doc)
    add_section_title(doc, 24, "Conclusion and References")
    add_body(
        doc,
        "The final project successfully advances the mid-semester TSFFM idea into a working AI-assisted depression screening prototype. It includes a reproducible feature pipeline, a two-stream temporal neural model, validation artifacts, backend inference, frontend interaction, and automated PDF reporting."
    )
    add_body(
        doc,
        "The current metrics show that the model is promising but not clinically ready. The best final evaluation framing is therefore balanced: the engineering integration is strong, the research concept is meaningful, and the results honestly identify generalization, recall, and validation as the main areas for future improvement."
    )
    add_matrix_table(
        doc,
        ["Reference", "Source"],
        [
            ("WHO depression fact sheet", "https://www.who.int/news-room/fact-sheets/detail/depression"),
            ("DAIC corpus", "Gratch et al., The Distress Analysis Interview Corpus, LREC 2014"),
            ("AVEC 2017 challenge", "Ringeval et al., Real-life Depression and Affect Recognition Workshop and Challenge"),
            ("LSTM", "Hochreiter and Schmidhuber, Neural Computation, 1997"),
            ("ResNet", "He et al., CVPR 2016"),
            ("Project source files", "README, ml/models.py, ml/train.py, ml/evaluate.py, backend services, frontend components"),
        ],
        (1.8, 4.5),
    )

    doc.core_properties.title = "TSFFM Depression Detection Final Evaluation Report"
    doc.core_properties.subject = "BTP final evaluation report"
    doc.core_properties.author = "Jitendra Choudhary and Kabir Singh Khair"
    doc.save(REPORT_PATH)
    return REPORT_PATH


PPT_BG = PptRGB(250, 251, 252)
PPT_INK = PptRGB(15, 37, 67)
PPT_TEAL = PptRGB(21, 112, 110)
PPT_COPPER = PptRGB(184, 115, 51)
PPT_MUTED = PptRGB(91, 106, 125)
PPT_LIGHT = PptRGB(232, 238, 245)
PPT_WARN = PptRGB(255, 244, 232)


def ppt_set_text(shape, text, size=18, color=PPT_INK, bold=False, align=None):
    shape.text = ""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def ppt_add_text(slide, text, x, y, w, h, size=18, color=PPT_INK, bold=False, align=None):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.margin_left = PptInches(0.05)
    tf.margin_right = PptInches(0.05)
    tf.margin_top = PptInches(0.02)
    tf.margin_bottom = PptInches(0.02)
    ppt_set_text(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def ppt_add_title(slide, kicker, title):
    ppt_add_text(slide, kicker.upper(), 0.55, 0.32, 3.2, 0.25, size=8.5, color=PPT_COPPER, bold=True)
    ppt_add_text(slide, title, 0.55, 0.58, 11.5, 0.65, size=27, color=PPT_INK, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.55), PptInches(1.27), PptInches(0.9), PptInches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = PPT_TEAL
    line.line.fill.background()


def ppt_add_footer(slide, idx):
    ppt_add_text(slide, "TSFFM Depression Detection | Final Evaluation", 0.55, 7.08, 5.5, 0.2, size=7.5, color=PPT_MUTED)
    ppt_add_text(slide, f"{idx:02d}", 12.45, 7.05, 0.35, 0.22, size=8, color=PPT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


def ppt_add_box(slide, x, y, w, h, title, body="", fill=PPT_LIGHT, accent=PPT_TEAL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = PptPt(1.0)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.14)
    tf.margin_right = PptInches(0.14)
    tf.margin_top = PptInches(0.09)
    tf.margin_bottom = PptInches(0.08)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = PptPt(13)
    r.font.bold = True
    r.font.color.rgb = accent
    if body:
        p2 = tf.add_paragraph()
        p2.space_before = PptPt(4)
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = "Aptos"
        r2.font.size = PptPt(10)
        r2.font.color.rgb = PPT_INK
    return box


def ppt_add_bullets(slide, items, x, y, w, h, size=14):
    shape = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = PptInches(0.05)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = PptPt(size)
        p.font.color.rgb = PPT_INK
        p.space_after = PptPt(6)
    return shape


def ppt_add_image(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w:
        kwargs["width"] = PptInches(w)
    if h:
        kwargs["height"] = PptInches(h)
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), **kwargs)
    return None


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PPT_BG
    return slide


def build_pptx(main, candidate, figures):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    # 1
    slide = blank_slide(prs)
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(7.5))
    block.fill.solid()
    block.fill.fore_color.rgb = PPT_INK
    block.line.fill.background()
    ppt_add_text(slide, "FINAL EVALUATION", 0.65, 0.55, 4, 0.3, size=10, color=PPT_COPPER, bold=True)
    ppt_add_text(slide, "Depression Detection using Two-Stream Feature Fusion Model", 0.65, 1.05, 9.8, 1.35, size=34, color=PptRGB(255, 255, 255), bold=True)
    ppt_add_text(slide, "TSFFM-BiLSTM with MediaPipe visual features, FastAPI inference, React dashboard, and PDF reporting", 0.70, 2.55, 8.9, 0.6, size=16, color=PptRGB(220, 230, 238))
    ppt_add_text(slide, TEAM, 0.70, 5.55, 7.5, 0.35, size=13, color=PptRGB(255, 255, 255), bold=True)
    ppt_add_text(slide, f"Supervisor: {SUPERVISOR} | {DEPARTMENT}", 0.70, 5.95, 8.4, 0.3, size=11, color=PptRGB(210, 220, 230))
    ppt_add_box(slide, 9.7, 4.75, 2.65, 1.25, "Prototype status", "End-to-end system implemented and evaluated", fill=PptRGB(244, 248, 248), accent=PPT_TEAL)

    # 2
    slide = blank_slide(prs)
    ppt_add_title(slide, "problem", "Why depression screening needs scalable assistance")
    ppt_add_box(slide, 0.65, 1.65, 3.2, 1.3, "332M", "people worldwide estimated to have depression by WHO", fill=PptRGB(244, 248, 248))
    ppt_add_box(slide, 4.2, 1.65, 3.2, 1.3, "Access gap", "Many cases remain untreated because care is limited or delayed", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_box(slide, 7.75, 1.65, 3.9, 1.3, "Behavior signal", "Interview video contains facial affect, posture, and movement cues", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_bullets(slide, ["Clinical interviews remain essential.", "AI can support early screening by surfacing objective visual cues.", "The safe role is assistive triage, not autonomous diagnosis."], 0.75, 3.55, 10.9, 1.6, size=17)
    ppt_add_footer(slide, 2)

    # 3
    slide = blank_slide(prs)
    ppt_add_title(slide, "progress", "Mid-sem idea evolved into a full final prototype")
    ppt_add_box(slide, 0.8, 1.7, 3.2, 3.6, "Mid-sem state", "Conceptual TSFFM deck\nFrame-level two-stream idea\nSingle demo probability\nCPU-bound early training\nFuture work: LSTM and deployment", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_box(slide, 5.0, 1.7, 3.2, 3.6, "Final model", "Landmark-based streams\n360-frame sequences\nTemporal Conv1D\nBiLSTM + attention\nTrained weights and history", fill=PptRGB(244, 248, 248), accent=PPT_TEAL)
    ppt_add_box(slide, 9.2, 1.7, 3.2, 3.6, "Final system", "FastAPI backend\nReact/Vite dashboard\nPDF report generation\nValidation plots\nEnd-to-end workflow", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_footer(slide, 3)

    # 4
    slide = blank_slide(prs)
    ppt_add_title(slide, "system", "Final inference pipeline is complete end to end")
    ppt_add_image(slide, figures["pipeline"], 0.65, 1.55, w=12.0)
    ppt_add_bullets(slide, ["Samples video at 5 FPS and standardizes to 360 frames.", "Extracts 68 facial landmarks and shoulder pose cues.", "Returns prediction, confidence, detection quality, and PDF report URL."], 0.9, 4.85, 11.4, 1.1, size=15)
    ppt_add_footer(slide, 4)

    # 5
    slide = blank_slide(prs)
    ppt_add_title(slide, "data", "Dataset design uses segmented clinical interview features")
    ppt_add_box(slide, 0.8, 1.7, 3.2, 1.6, "Input format", "Participant directories with face-keypoint, pose, and binary-label NumPy files")
    ppt_add_box(slide, 4.15, 1.7, 3.2, 1.6, "Validation set", "352 segments\n211 not depressed\n141 depressed", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_box(slide, 7.5, 1.7, 3.9, 1.6, "Label target", "0 = not depressed\n1 = depressed\nBinary PHQ-style screening setup", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_bullets(slide, ["Training uses class weights, focal loss, and balanced sampling.", "Temporal jittering and Gaussian landmark noise help regularization.", "Participant-level diversity remains the major data challenge."], 0.9, 4.15, 11.0, 1.3, size=16)
    ppt_add_footer(slide, 5)

    # 6
    slide = blank_slide(prs)
    ppt_add_title(slide, "features", "Preprocessing converts video into normalized temporal landmarks")
    rows = [
        ("Face", "68 landmarks x 4 values", "272-D per frame"),
        ("Pose", "2 shoulder points x 4 values", "8-D per frame"),
        ("Sequence", "Pad/truncate to 360 frames", "Fixed model input"),
    ]
    x0, y0 = 0.9, 1.75
    for i, (a, b, c) in enumerate(rows):
        ppt_add_box(slide, x0, y0 + i * 1.25, 2.35, 0.8, a, b, fill=PptRGB(244, 248, 248), accent=PPT_TEAL)
        ppt_add_box(slide, x0 + 3.2, y0 + i * 1.25, 3.1, 0.8, "Model input", c, fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_bullets(slide, ["Each frame is centered and scale-normalized.", "Missing detections reuse last available features where possible.", "Detection rates are returned to the frontend for quality visibility."], 7.8, 1.85, 4.6, 2.5, size=15)
    ppt_add_footer(slide, 6)

    # 7
    slide = blank_slide(prs)
    ppt_add_title(slide, "model", "Final architecture adds temporal learning and attention")
    ppt_add_image(slide, figures["model"], 0.75, 1.45, w=11.8)
    ppt_add_bullets(slide, ["Face stream: 272 -> 128", "Body stream: 8 -> 32", "Fusion: 160-D sequence", "Temporal Conv1D + BiLSTM + attention -> 2 logits"], 0.9, 5.25, 11.4, 0.7, size=14)
    ppt_add_footer(slide, 7)

    # 8
    slide = blank_slide(prs)
    ppt_add_title(slide, "training", "Training was designed for imbalance and small-data risk")
    items = [
        ("Focal loss", "Down-weights easy examples and focuses on harder cases"),
        ("Class weights", "Compensates for depressed / not-depressed imbalance"),
        ("Balanced sampler", "Improves exposure to minority-class segments"),
        ("Regularization", "Dropout, weight decay, gradient clipping, LR scheduling"),
        ("Checkpointing", "Best validation metric saved for final evaluation"),
    ]
    for i, (title, body) in enumerate(items):
        ppt_add_box(slide, 0.85 + (i % 3) * 4.05, 1.7 + (i // 3) * 1.65, 3.55, 1.1, title, body, fill=PptRGB(244, 248, 248) if i % 2 == 0 else PPT_LIGHT)
    ppt_add_footer(slide, 8)

    # 9
    slide = blank_slide(prs)
    ppt_add_title(slide, "evaluation", "Final evaluation uses both performance and error-shape metrics")
    ppt_add_box(slide, 0.8, 1.7, 2.4, 1.3, "Accuracy", "Overall correctness")
    ppt_add_box(slide, 3.45, 1.7, 2.4, 1.3, "Precision", "Reliability of positive screenings", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_box(slide, 6.1, 1.7, 2.4, 1.3, "Recall", "How many depressed cases are detected", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_box(slide, 8.75, 1.7, 2.4, 1.3, "AUC-ROC", "Probability ranking quality")
    ppt_add_bullets(slide, ["For screening, recall is critical because false negatives are high-risk.", "For practical use, recall must be balanced against false-positive burden.", "The current prototype is evaluated as research software, not a clinical device."], 0.9, 4.05, 11.2, 1.4, size=16)
    ppt_add_footer(slide, 9)

    # 10
    slide = blank_slide(prs)
    ppt_add_title(slide, "results", "Selected checkpoint is useful, but sensitivity needs work")
    ppt_add_image(slide, figures["metrics"], 0.75, 1.45, w=7.3)
    ppt_add_box(slide, 8.6, 1.6, 3.2, 1.1, "Accuracy", f"{main['accuracy']:.2f}%", fill=PptRGB(244, 248, 248))
    ppt_add_box(slide, 8.6, 2.95, 3.2, 1.1, "F1-score", f"{main['f1']:.4f}", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_box(slide, 8.6, 4.3, 3.2, 1.1, "AUC-ROC", f"{main['auc']:.4f}", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_footer(slide, 10)

    # 11
    slide = blank_slide(prs)
    ppt_add_title(slide, "errors", "Confusion matrix shows the screening trade-off")
    ppt_add_image(slide, RESULTS / "confusion_matrix.png", 0.9, 1.55, w=4.7)
    ppt_add_image(slide, RESULTS / "roc_curve.png", 6.15, 1.55, w=4.7)
    ppt_add_bullets(slide, ["Selected model: better accuracy and AUC.", "Recall-oriented candidate: stronger depression recall but lower AUC.", "Final threshold should be calibrated with clinical input."], 1.0, 5.75, 10.8, 0.75, size=13.5)
    ppt_add_footer(slide, 11)

    # 12
    slide = blank_slide(prs)
    ppt_add_title(slide, "learning", "Training curves reveal overfitting pressure")
    ppt_add_image(slide, RESULTS / "accuracy_history.png", 0.85, 1.55, w=5.25)
    ppt_add_image(slide, RESULTS / "loss_history.png", 6.6, 1.55, w=5.25)
    ppt_add_box(slide, 1.2, 5.65, 10.7, 0.65, "Interpretation", "Training improves faster than validation, so future work should prioritize data diversity, regularization, and participant-level split checks.", fill=PptRGB(244, 248, 248))
    ppt_add_footer(slide, 12)

    # 13
    slide = blank_slide(prs)
    ppt_add_title(slide, "implementation", "The final project is productized as a usable web workflow")
    ppt_add_image(slide, figures["workflow"], 0.7, 1.4, w=11.6)
    ppt_add_bullets(slide, ["Frontend: upload, metadata, progress, result card, probability chart.", "Backend: validation, extraction, model singleton, PDF generation, cleanup.", "Report: patient/session metadata, prediction, probabilities, disclaimer."], 0.9, 5.35, 11.3, 0.95, size=14)
    ppt_add_footer(slide, 13)

    # 14
    slide = blank_slide(prs)
    ppt_add_title(slide, "safety", "Ethics and privacy define the deployment boundary")
    ppt_add_box(slide, 0.9, 1.65, 3.4, 1.55, "Not a diagnosis", "The system is a screening assistant and requires professional review.", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_box(slide, 4.75, 1.65, 3.4, 1.55, "Sensitive data", "Interview videos are biometric and behavioral health information.", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_box(slide, 8.6, 1.65, 3.4, 1.55, "Bias risk", "Performance must be audited across demographics and recording conditions.")
    ppt_add_bullets(slide, ["Use consent, local processing, deletion policies, and access controls.", "Do not use outputs for high-stakes decisions.", "Clinical validation is mandatory before real-world evaluation."], 1.0, 4.2, 10.8, 1.2, size=16)
    ppt_add_footer(slide, 14)

    # 15
    slide = blank_slide(prs)
    ppt_add_title(slide, "roadmap", "Next steps focus on stronger evidence and generalization")
    roadmap = [
        ("1", "Full pose stream", "Use richer upper-body landmarks and motion descriptors"),
        ("2", "Audio + text", "Add voice and transcript features for true multimodal fusion"),
        ("3", "Threshold calibration", "Tune screening sensitivity with clinical guidance"),
        ("4", "External validation", "Test on larger and more diverse datasets"),
    ]
    for i, (num, title, body) in enumerate(roadmap):
        ppt_add_box(slide, 0.9 + (i % 2) * 5.75, 1.75 + (i // 2) * 1.75, 4.9, 1.25, f"{num}. {title}", body, fill=PptRGB(244, 248, 248) if i % 2 == 0 else PPT_LIGHT, accent=PPT_TEAL if i % 2 == 0 else PPT_INK)
    ppt_add_footer(slide, 15)

    # 16
    slide = blank_slide(prs)
    ppt_add_title(slide, "conclusion", "Final evaluation: strong prototype, honest research frontier")
    ppt_add_box(slide, 0.9, 1.65, 3.5, 2.4, "What is complete", "TSFFM model\nTemporal sequence learning\nValidation artifacts\nFastAPI inference\nReact dashboard\nPDF reports", fill=PptRGB(244, 248, 248))
    ppt_add_box(slide, 4.9, 1.65, 3.5, 2.4, "What results show", "Useful learned signal\nAUC above chance\nSensitivity still weak\nOverfitting pressure\nNeed calibration", fill=PPT_LIGHT, accent=PPT_INK)
    ppt_add_box(slide, 8.9, 1.65, 3.5, 2.4, "Final claim", "A meaningful AI-assisted screening prototype, ready for demo and further research, not clinical deployment.", fill=PptRGB(255, 244, 232), accent=PPT_COPPER)
    ppt_add_text(slide, "Thank You", 0.9, 5.55, 4.0, 0.5, size=28, color=PPT_INK, bold=True)
    ppt_add_text(slide, TEAM, 0.95, 6.12, 6.2, 0.3, size=12, color=PPT_MUTED)
    ppt_add_footer(slide, 16)

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main():
    ensure_dirs()
    main_metrics, candidate_metrics, history = read_metrics()
    figures = generate_figures(main_metrics, candidate_metrics)
    report = build_report(main_metrics, candidate_metrics, history, figures)
    ppt = build_pptx(main_metrics, candidate_metrics, figures)
    print(f"Report: {report}")
    print(f"PPTX: {ppt}")


if __name__ == "__main__":
    main()
